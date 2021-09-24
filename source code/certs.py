from pptx import Presentation
import tkinter as tk
from tkinter import filedialog
import pandas as pd
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
from pptx.enum.dml import MSO_THEME_COLOR
import os
import win32com.client
import sys

#all folders and directories used
cwd = os.getcwd()

startfolder=os.path.join(os.environ['USERPROFILE'], 'Desktop', 'Certificates', 'PPTs')
if not os.path.exists(startfolder):
    os.makedirs(os.path.join(os.environ['USERPROFILE'], 'Desktop', 'Certificates', 'PPTs'))

destinyfolder=os.path.join(os.environ['USERPROFILE'], 'Desktop', 'Certificates', 'PDFs')
if not os.path.exists(destinyfolder):
    os.makedirs(os.path.join(os.environ['USERPROFILE'], 'Desktop', 'Certificates', 'PDFs'))


def wholeprogram():
    usercom=int(input("Hello user. Press 1 to start or 0 to quit."))
    if usercom==1:
        #take data from excel sheet
        print("Hello user. Please choose an excel file. Press enter to choose.")
        input()

        root = tk.Tk()
        root.withdraw()
        file_path = filedialog.askopenfilename(filetypes=[("Excel files","*.xlsx;*.xls")])
        df=pd.read_excel(file_path)
        list_cols=tuple(df)


        #open ppt
        print("Now please choose a ppt file. Press enter to choose.")
        input()

        root = tk.Tk()
        root.withdraw()
        file_path1 = filedialog.askopenfilename(filetypes=[("Powerpoint files","*.pptx;*.ppt")])
        prs = Presentation(file_path1)
        
        h = input('Enter hex color code. Eg- #FFFFFF').lstrip('#')
        rgbvals=tuple(int(h[i:i+2], 16) for i in (0, 2, 4))
        rval=rgbvals[0]
        gval=rgbvals[1]
        bval=rgbvals[2]

        #change text
        for slide in prs.slides:
            for i in range(len(df.index)):
                for item in list_cols:
                    for shape in slide.shapes:
                        if (not shape.has_text_frame) or (shape.text in list_cols): 
                            continue
                        else:
                            shape.text=item
                            break
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue 
                    for column in df:
                        if shape.text!=column:
                            continue
                        val1=str(df[column].values[i])
                        shape.text=val1.upper()
                        text_frame = shape.text_frame
                        text_frame.fit_text(u'Montserrat-Medium', max_size=40, bold=False, italic=False, font_file=os.path.join(os.environ['USERPROFILE'], 'AppData', 'Local', 'Microsoft', 'Windows', 'Fonts', 'Montserrat-Medium.otf'))
                        p = text_frame.paragraphs[0]
                        p.alignment = PP_ALIGN.CENTER
                        font=p.font
                        font.color.rgb = RGBColor(rval, gval, bval)
                pptx_name=str(df['NAME'].values[i])
                prs.save(os.path.join(os.environ['USERPROFILE'], 'Desktop', 'Certificates', 'PPTs', pptx_name.upper()+'.pptx'))
            
        #%% Get console arguments
        input_folder_path = startfolder
        output_folder_path = destinyfolder

        #%% Get files in input folder
        input_file_paths = os.listdir(input_folder_path)

        #%% Convert each file
        for input_file_name in input_file_paths:

            # Skip if file does not contain a power point extension
            if not input_file_name.lower().endswith((".ppt", ".pptx")):
                continue
            
            # Create input file path
            input_file_path = os.path.join(input_folder_path, input_file_name)
                
            # Create powerpoint application object
            powerpoint = win32com.client.DispatchEx("Powerpoint.Application")
            
            # Set visibility to minimize
            powerpoint.Visible = 1
            
            # Open the powerpoint slides
            slides = powerpoint.Presentations.Open(input_file_path)
            
            # Get base file name
            file_name = os.path.splitext(input_file_name)[0]
            
            # Create output file path
            output_file_path = os.path.join(output_folder_path, file_name + ".pdf")
            
            # Save as PDF (formatType = 32)
            slides.SaveAs(output_file_path, 32)
            
            # Close the slide deck
            slides.Close()

        wholeprogram()
    
    else:
        pass

wholeprogram()