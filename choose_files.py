from pptx import Presentation
import tkinter as tk
from tkinter import filedialog
import pandas as pd
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from pathlib import Path
from pptx.enum.dml import MSO_THEME_COLOR
import os
import win32com.client
import sys
import eel
import matplotlib

eel.init('web')

cwd = os.getcwd()

@eel.expose
def choose_excel():
    root = tk.Tk()
    root.withdraw()
    file_path_excel = filedialog.askopenfilename(filetypes=[("Excel files","*.xlsx;*.xls")])
    df=pd.read_excel(file_path_excel)
    df_html=df.to_html()
    return file_path_excel, df_html

@eel.expose
def choose_ppt():
    root = tk.Tk()
    root.withdraw()
    file_path_ppt = filedialog.askopenfilename(filetypes=[("Powerpoint files","*.pptx;*.ppt")])
    return file_path_ppt

@eel.expose
def choose_save():
    root = tk.Tk()
    root.withdraw()
    file_path_save = filedialog.askdirectory()
    return file_path_save


def make_pdfs(file_path_excel, file_path_ppt, file_path_save, color):
    #all folders and directories used
    if (file_path_excel!="") and (file_path_ppt!="") and (file_path_save!=""):
        startfolder=os.path.join(file_path_save, 'Certificates', 'PPTs')
        if not os.path.exists(startfolder):
            os.makedirs(startfolder)

        destinyfolder=os.path.join(file_path_save, 'Certificates', 'PDFs')
        if not os.path.exists(destinyfolder):
            os.makedirs(destinyfolder)


        df=pd.read_excel(file_path_excel)
        list_cols=tuple(df)

        #change text
        prs = Presentation(file_path_ppt)
        for slide in prs.slides:
            for i in range(len(df.index)):
                for item in list_cols:
                    for shape in slide.shapes:
                        if (not shape.has_text_frame) or (shape.text in list_cols): 
                            continue
                        else:
                            shape.text=item
                            break
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue 
                    for column in df:
                        if shape.text!=column:
                            continue
                        val1=str(df[column].values[i])
                        shape.text=val1.upper()
                        text_frame = shape.text_frame
                        text_frame.fit_text(u'Montserrat-Medium', max_size=40, bold=False, italic=False, font_file=os.path.join(os.environ['USERPROFILE'], 'AppData', 'Local', 'Microsoft', 'Windows', 'Fonts', 'Montserrat-Medium.otf'))
                        p = text_frame.paragraphs[0]
                        p.alignment = PP_ALIGN.CENTER
                        font=p.font
                        font.color=color
                pptx_name=str(df['NAME'].values[i])
                prs.save(os.path.join(startfolder, pptx_name.upper()+'.pptx'))
            
        #%% Get console arguments
        input_folder_path = startfolder
        output_folder_path = destinyfolder

        #%% Get files in input folder
        input_file_paths = os.listdir(input_folder_path)

        #%% Convert each file
        for input_file_name in input_file_paths:

            # Skip if file does not contain a power point extension
            if not input_file_name.lower().endswith((".ppt", ".pptx")):
                continue
            
            # Create input file path
            input_file_path = os.path.join(input_folder_path, input_file_name)
                
            # Create powerpoint application object
            powerpoint = win32com.client.DispatchEx("Powerpoint.Application")
            
            # Set visibility to minimize
            powerpoint.Visible = 1
            
            # Open the powerpoint slides
            slides = powerpoint.Presentations.Open(input_file_path)
            
            # Get base file name
            file_name = os.path.splitext(input_file_name)[0]
            
            # Create output file path
            output_file_path = os.path.join(output_folder_path, file_name + ".pdf")
            
            # Save as PDF (formatType = 32)
            slides.SaveAs(output_file_path, 32)
            
            # Close the slide deck
            slides.Close()
        return "green", "done"
    else:
        elems={file_path_excel:"Excel file not selected", file_path_ppt:"Ppt file not selected", file_path_save:"Output folder not selected"}
        errs=[]
        for key, value in elems.items():
            if key=="":
                if "red" not in errs:
                    errs.append("red")
                errs.append(value)
        return errs

eel.start('gui.html', size=(1000, 750), mode='default')